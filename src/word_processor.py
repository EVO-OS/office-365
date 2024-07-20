# Make sure to install the required dependencies:
# pip install python-docx openpyxl python-pptx

import gi
import os
import sys
import docx
from openpyxl import load_workbook, Workbook
from pptx import Presentation
from ai_integration import generate_text, summarize_text, translate_text, analyze_sentiment, answer_question

gi.require_version('Gtk', '3.0')
from gi.repository import Gtk, Gdk, GLib

print("Python version:", sys.version)
print("GTK version:", gi.version_info)

def gtk_main_thread(func):
    def wrapper(*args, **kwargs):
        GLib.idle_add(func, *args, **kwargs)
    return wrapper

def initialize_gtk():
    print("Attempting to initialize GTK...")
    success, argv = Gtk.init_check(sys.argv)
    if not success:
        print("Unable to initialize GTK.", file=sys.stderr)
        print("Gtk.init_check() returned:", success)
        print("Modified argv:", argv)
        sys.exit(1)
    print("GTK initialized successfully")

print("About to call initialize_gtk()")
initialize_gtk()
print("initialize_gtk() completed")

initialize_gtk()
print(f"GTK version: {Gtk.get_major_version()}.{Gtk.get_minor_version()}.{Gtk.get_micro_version()}")

class TextEditor:
    def __init__(self):
        self.text_buffer = ""
        self.current_file = None
        self.formatting = {}

    def load_file(self, filepath):
        with open(filepath, 'r') as file:
            self.text_buffer = file.read()
        self.current_file = filepath
        self.formatting = {}

    def save_file(self, filepath):
        with open(filepath, 'w') as file:
            file.write(self.text_buffer)
        self.current_file = filepath

    def format_text(self, style, start, end):
        if start < 0 or end > len(self.text_buffer) or start >= end:
            raise ValueError("Invalid range")

        for i in range(start, end):
            if i not in self.formatting:
                self.formatting[i] = set()
            self.formatting[i].add(style)

    def insert_image(self, image_path, position):
        if position < 0 or position > len(self.text_buffer):
            raise ValueError("Invalid position")

        image_marker = f"[IMAGE:{image_path}]"
        self.text_buffer = self.text_buffer[:position] + image_marker + self.text_buffer[position:]

    def get_formatted_text(self):
        formatted_text = ""
        for i, char in enumerate(self.text_buffer):
            if i in self.formatting:
                styles = self.formatting[i]
                formatted_char = char
                for style in styles:
                    if style == 'bold':
                        formatted_char = f"<b>{formatted_char}</b>"
                    elif style == 'italic':
                        formatted_char = f"<i>{formatted_char}</i>"
                    elif style == 'underline':
                        formatted_char = f"<u>{formatted_char}</u>"
                formatted_text += formatted_char
            else:
                formatted_text += char
        return formatted_text

def import_document(file_path):
    """
    Import a document from Google Office or Microsoft Office 365 format.
    """
    import docx
    from openpyxl import load_workbook
    from pptx import Presentation

    file_extension = file_path.split('.')[-1].lower()

    if file_extension in ['docx', 'xlsx', 'pptx']:
        if file_extension == 'docx':
            doc = docx.Document(file_path)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        elif file_extension == 'xlsx':
            wb = load_workbook(file_path)
            sheet = wb.active
            return '\n'.join(['\t'.join([str(cell.value) for cell in row]) for row in sheet.rows])
        elif file_extension == 'pptx':
            prs = Presentation(file_path)
            return '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")

def export_document(content, file_path):
    """
    Export a document to Google Office or Microsoft Office 365 format.
    """
    import docx
    from openpyxl import Workbook
    from pptx import Presentation

    file_extension = file_path.split('.')[-1].lower()

    if file_extension in ['docx', 'xlsx', 'pptx']:
        if file_extension == 'docx':
            doc = docx.Document()
            for paragraph in content.split('\n'):
                doc.add_paragraph(paragraph)
            doc.save(file_path)
        elif file_extension == 'xlsx':
            wb = Workbook()
            sheet = wb.active
            for row, line in enumerate(content.split('\n'), start=1):
                for col, value in enumerate(line.split('\t'), start=1):
                    sheet.cell(row=row, column=col, value=value)
            wb.save(file_path)
        elif file_extension == 'pptx':
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = "Exported Document"
            body = slide.shapes.placeholders[1]
            body.text = content
            prs.save(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")

class WordProcessorUI:
    def __init__(self):
        print("WordProcessorUI.__init__() started")
        print("Current working directory:", os.getcwd())
        print("DISPLAY environment variable:", os.environ.get('DISPLAY'))
        print("Initializing WordProcessorUI")
        try:
            self.builder = self.load_glade_file()
            self.connect_signals()
            css_provider = self.create_css_provider()
            self.apply_css_provider(css_provider)
            self.initialize_main_window()
            self.initialize_text_view()
            self.initialize_toolbar()
            self.editor = TextEditor()

            # Add import and export buttons
            self.import_button = self.builder.get_object("import_button")
            self.export_button = self.builder.get_object("export_button")

            # Add AI feature buttons
            self.ai_generate_button = self.builder.get_object("ai_generate_button")
            self.ai_summarize_button = self.builder.get_object("ai_summarize_button")
            self.ai_translate_button = self.builder.get_object("ai_translate_button")
            self.ai_sentiment_button = self.builder.get_object("ai_sentiment_button")
            self.ai_qa_button = self.builder.get_object("ai_qa_button")

            # Connect signals for import, export, and AI feature buttons
            self.import_button.connect("clicked", self.on_import_clicked)
            self.export_button.connect("clicked", self.on_export_clicked)
            self.ai_generate_button.connect("clicked", self.on_ai_generate_clicked)
            self.ai_summarize_button.connect("clicked", self.on_ai_summarize_clicked)
            self.ai_translate_button.connect("clicked", self.on_ai_translate_clicked)
            self.ai_sentiment_button.connect("clicked", self.on_ai_sentiment_clicked)
            self.ai_qa_button.connect("clicked", self.on_ai_qa_clicked)
        except Exception as e:
            print(f"Error during initialization: {e}")
            raise

    def load_glade_file(self):
        print("Loading Glade file")
        print(f"Current working directory: {os.getcwd()}")
        glade_file_path = os.path.abspath('word_processor_ui.glade')
        print(f"Glade file path: {glade_file_path}")

        builder = Gtk.Builder()
        try:
            print("Attempting to load Glade file...")
            builder.add_from_file(glade_file_path)
            print("Glade file loaded successfully")
            return builder
        except GLib.Error as e:
            print(f"Error loading Glade file: {e}")
            print(f"Error details: {e.args}")
            raise

    def connect_signals(self):
        print("Connecting signals")
        self.builder.connect_signals(self)
        print("Signals connected")

    def create_css_provider(self):
        print("Creating CSS provider")
        css_provider = Gtk.CssProvider()
        css_file_path = os.path.abspath('word_processor_styles.css')
        print(f"Attempting to load CSS file: {css_file_path}")

        try:
            css_provider.load_from_path(css_file_path)
            print("CSS file loaded successfully")
        except GLib.Error as e:
            print(f"Error loading CSS file: {e}")
            print("Falling back to inline CSS")
            self.load_fallback_css(css_provider)

        return css_provider

    def load_fallback_css(self, css_provider):
        css_rules = [
            (".libreoffice-button", """
                border-radius: 2px;
                border-width: 1px;
                border-color: #d0d0d0;
                background-color: #f0f0f0;
                color: #333333;
                padding: 4px 8px;
            """),
            (".sidebar-section", """
                background-color: #f5f5f5;
                border-bottom: 1px solid #d0d0d0;
                padding: 8px;
            """),
            (".suggested-action", """
                background-color: #4a90d9;
                color: #ffffff;
            """)
        ]

        for selector, rules in css_rules:
            print(f"Applying CSS rule: {selector}")
            try:
                css_provider.load_from_data(f"{selector} {{ {rules} }}".encode())
                print(f"CSS rule applied: {selector}")
            except GLib.Error as e:
                print(f"Error applying CSS rule {selector}: {e}")

        print("Fallback CSS loading complete")

    def initialize_main_window(self):
        print("Initializing main window")
        self.window = self.builder.get_object("main_window")
        if not self.window:
            raise ValueError("Main window not found in Glade file")
        print(f"Main window object type: {type(self.window)}")
        self.window.connect("destroy", Gtk.main_quit)
        print("Destroy signal connected successfully")

    def initialize_text_view(self):
        print("Initializing text view")
        self.text_view = self.builder.get_object("text_view")
        if not self.text_view:
            raise ValueError("Text view not found in Glade file")
        print(f"Text view object type: {type(self.text_view)}")
        self.text_buffer = self.text_view.get_buffer()
        print(f"Text buffer type: {type(self.text_buffer)}")

    def initialize_toolbar(self):
        print("Initializing toolbar")
        toolbar = self.builder.get_object("toolbar")
        print(f"Toolbar object: {toolbar}")

        button_ids = ["new_button", "open_button", "save_button", "bold_button", "italic_button", "underline_button",
                      "ai_generate_button", "ai_summarize_button", "ai_translate_button", "ai_sentiment_button", "ai_qa_button"]
        for button_id in button_ids:
            print(f"Adding tool button: {button_id}")
            button = self.builder.get_object(button_id)
            if button:
                toolbar.insert(button, -1)
            else:
                print(f"Warning: Button {button_id} not found in Glade file")

    def apply_css_provider(self, css_provider):
        try:
            print("Attempting to add CSS provider for screen")
            screen = Gdk.Screen.get_default()
            Gtk.StyleContext.add_provider_for_screen(
                screen,
                css_provider,
                Gtk.STYLE_PROVIDER_PRIORITY_APPLICATION
            )
            print("CSS provider added successfully")
            return True
        except GLib.Error as e:
            print(f"Error applying CSS provider: {e}")
            return False

    def run(self):
        print("Starting WordProcessorUI")
        print("About to show main window")
        self.window.show_all()
        print("Main window shown")
        print("About to enter Gtk main loop")
        try:
            Gtk.main()
        except KeyboardInterrupt:
            print("Application terminated by user")
        except Exception as e:
            print(f"Error in Gtk main loop: {e}")
        finally:
            print("Exited Gtk main loop")
            Gtk.main_quit()

    def on_import_clicked(self, widget):
        dialog = Gtk.FileChooserDialog(
            title="Import Document", parent=self.window, action=Gtk.FileChooserAction.OPEN
        )
        dialog.add_buttons(
            Gtk.STOCK_CANCEL, Gtk.ResponseType.CANCEL, Gtk.STOCK_OPEN, Gtk.ResponseType.OK
        )

        filter_docx = Gtk.FileFilter()
        filter_docx.set_name("Office Documents")
        filter_docx.add_pattern("*.docx")
        filter_docx.add_pattern("*.xlsx")
        filter_docx.add_pattern("*.pptx")
        dialog.add_filter(filter_docx)

        response = dialog.run()
        if response == Gtk.ResponseType.OK:
            file_path = dialog.get_filename()
            content = import_document(file_path)
            self.text_buffer.set_text(content)

        dialog.destroy()

    def on_export_clicked(self, widget):
        dialog = Gtk.FileChooserDialog(
            title="Export Document", parent=self.window, action=Gtk.FileChooserAction.SAVE
        )
        dialog.add_buttons(
            Gtk.STOCK_CANCEL, Gtk.ResponseType.CANCEL, Gtk.STOCK_SAVE, Gtk.ResponseType.OK
        )

        filter_docx = Gtk.FileFilter()
        filter_docx.set_name("Office Documents")
        filter_docx.add_pattern("*.docx")
        filter_docx.add_pattern("*.xlsx")
        filter_docx.add_pattern("*.pptx")
        dialog.add_filter(filter_docx)

        response = dialog.run()
        if response == Gtk.ResponseType.OK:
            file_path = dialog.get_filename()
            start_iter = self.text_buffer.get_start_iter()
            end_iter = self.text_buffer.get_end_iter()
            content = self.text_buffer.get_text(start_iter, end_iter, True)
            export_document(content, file_path)

        dialog.destroy()

    def on_ai_generate_clicked(self, widget):
        prompt = self.get_selected_text() or "Generate text about:"
        generated_text = generate_text(prompt)
        self.insert_text_at_cursor(generated_text)

    def on_ai_summarize_clicked(self, widget):
        text = self.get_selected_text() or self.get_all_text()
        summary = summarize_text(text)
        self.insert_text_at_cursor("\n\nSummary:\n" + summary)

    def on_ai_translate_clicked(self, widget):
        text = self.get_selected_text() or self.get_all_text()
        translated_text = translate_text(text, "en")  # Translate to English
        self.insert_text_at_cursor("\n\nTranslation:\n" + translated_text)

    def on_ai_sentiment_clicked(self, widget):
        text = self.get_selected_text() or self.get_all_text()
        sentiment = analyze_sentiment(text)
        self.insert_text_at_cursor("\n\nSentiment: " + sentiment)

    def on_ai_qa_clicked(self, widget):
        context = self.get_all_text()
        question = self.show_input_dialog("Enter your question:")
        if question:
            answer = answer_question(context, question)
            self.insert_text_at_cursor("\n\nQ: " + question + "\nA: " + answer)

    def get_selected_text(self):
        bounds = self.text_buffer.get_selection_bounds()
        if bounds:
            return self.text_buffer.get_text(bounds[0], bounds[1], True)
        return ""

    def get_all_text(self):
        start_iter = self.text_buffer.get_start_iter()
        end_iter = self.text_buffer.get_end_iter()
        return self.text_buffer.get_text(start_iter, end_iter, True)

    def insert_text_at_cursor(self, text):
        self.text_buffer.insert_at_cursor(text)

    def show_input_dialog(self, message):
        dialog = Gtk.MessageDialog(
            transient_for=self.window,
            flags=0,
            message_type=Gtk.MessageType.QUESTION,
            buttons=Gtk.ButtonsType.OK_CANCEL,
            text=message,
        )
        entry = Gtk.Entry()
        entry.set_activates_default(True)
        dialog.get_content_area().add(entry)
        dialog.set_default_response(Gtk.ResponseType.OK)
        dialog.show_all()
        response = dialog.run()
        text = entry.get_text()
        dialog.destroy()
        if response == Gtk.ResponseType.OK:
            return text
        return None

if __name__ == "__main__":
    print("Main block started")
    try:
        print("About to create WordProcessorUI instance")
        app = WordProcessorUI()
        print("WordProcessorUI instance created successfully")
        print("About to call app.run()")
        app.run()
    except Exception as e:
        print(f"Error initializing application: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

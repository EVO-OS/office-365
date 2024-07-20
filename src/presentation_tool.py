import gi
gi.require_version('Gtk', '3.0')
from gi.repository import Gtk, Gdk, GLib
import os
import sys
from llama3.model import Llama3Model
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt

class Presentation:
    def __init__(self):
        self.pptx = PptxPresentation()
        self.current_slide = None

    def new_slide(self):
        layout = self.pptx.slide_layouts[1]  # Using layout with title and content
        self.current_slide = self.pptx.slides.add_slide(layout)
        return self.current_slide

    def load_presentation(self, filepath):
        try:
            self.pptx = PptxPresentation(filepath)
            if self.pptx.slides:
                self.current_slide = self.pptx.slides[0]
            return True
        except Exception as e:
            print(f"Error loading presentation: {e}")
            return False

    def save_presentation(self, filepath):
        try:
            self.pptx.save(filepath)
            return True
        except Exception as e:
            print(f"Error saving presentation: {e}")
            return False

    def insert_media(self, media_path):
        if self.current_slide is not None:
            try:
                self.current_slide.shapes.add_picture(media_path, Inches(1), Inches(1))
                return True
            except Exception as e:
                print(f"Error inserting media: {e}")
        return False

    def update_slide_content(self, title, content):
        if self.current_slide is not None:
            title_shape = self.current_slide.shapes.title
            content_shape = self.current_slide.placeholders[1]

            if title_shape:
                title_shape.text = title
            if content_shape:
                content_shape.text = content

class PresentationUI(Gtk.Window):
    def __init__(self):
        super().__init__(title="Advanced Presentation Tool")
        self.set_default_size(1024, 768)

        # Initialize llama3 model
        self.llama_model = Llama3Model()

        self.presentation = Presentation()
        self.initialize_ui()

    def initialize_ui(self):
        # Main layout
        main_box = Gtk.Box(orientation=Gtk.Orientation.VERTICAL, spacing=6)
        self.add(main_box)

        # Toolbar
        toolbar = Gtk.Toolbar()
        main_box.pack_start(toolbar, False, False, 0)

        new_button = Gtk.ToolButton.new_from_stock(Gtk.STOCK_NEW)
        new_button.connect("clicked", self.on_new_clicked)
        toolbar.insert(new_button, -1)

        open_button = Gtk.ToolButton.new_from_stock(Gtk.STOCK_OPEN)
        open_button.connect("clicked", self.on_open_clicked)
        toolbar.insert(open_button, -1)

        save_button = Gtk.ToolButton.new_from_stock(Gtk.STOCK_SAVE)
        save_button.connect("clicked", self.on_save_clicked)
        toolbar.insert(save_button, -1)

        # Slide view
        self.slide_view = Gtk.TextView()
        main_box.pack_start(self.slide_view, True, True, 0)

        # AI features
        ai_box = Gtk.Box(orientation=Gtk.Orientation.HORIZONTAL, spacing=6)
        main_box.pack_start(ai_box, False, False, 0)

        generate_button = Gtk.Button(label="Generate Content")
        generate_button.connect("clicked", self.on_generate_clicked)
        ai_box.pack_start(generate_button, False, False, 0)

        research_button = Gtk.Button(label="Research Assistant")
        research_button.connect("clicked", self.on_research_clicked)
        ai_box.pack_start(research_button, False, False, 0)

    def on_new_clicked(self, widget):
        self.presentation = Presentation()
        self.presentation.new_slide()
        self.update_slide_view()

    def on_open_clicked(self, widget):
        dialog = Gtk.FileChooserDialog("Please choose a file", self,
            Gtk.FileChooserAction.OPEN,
            (Gtk.STOCK_CANCEL, Gtk.ResponseType.CANCEL,
             Gtk.STOCK_OPEN, Gtk.ResponseType.OK))

        filter_pptx = Gtk.FileFilter()
        filter_pptx.set_name("PowerPoint files")
        filter_pptx.add_pattern("*.pptx")
        dialog.add_filter(filter_pptx)

        response = dialog.run()
        if response == Gtk.ResponseType.OK:
            if self.presentation.load_presentation(dialog.get_filename()):
                self.update_slide_view()
            else:
                self.show_error_dialog("Failed to load presentation")
        dialog.destroy()

    def on_save_clicked(self, widget):
        dialog = Gtk.FileChooserDialog("Save file", self,
            Gtk.FileChooserAction.SAVE,
            (Gtk.STOCK_CANCEL, Gtk.ResponseType.CANCEL,
             Gtk.STOCK_SAVE, Gtk.ResponseType.OK))

        filter_pptx = Gtk.FileFilter()
        filter_pptx.set_name("PowerPoint files")
        filter_pptx.add_pattern("*.pptx")
        dialog.add_filter(filter_pptx)

        response = dialog.run()
        if response == Gtk.ResponseType.OK:
            filepath = dialog.get_filename()
            if not filepath.endswith('.pptx'):
                filepath += '.pptx'
            if not self.presentation.save_presentation(filepath):
                self.show_error_dialog("Failed to save presentation")
        dialog.destroy()

    def update_slide_view(self):
        if self.presentation.current_slide:
            title = self.presentation.current_slide.shapes.title.text if self.presentation.current_slide.shapes.title else ""
            content = self.presentation.current_slide.placeholders[1].text if len(self.presentation.current_slide.placeholders) > 1 else ""
            self.slide_view.get_buffer().set_text(f"{title}\n\n{content}")

    def on_generate_clicked(self, widget):
        prompt = "Generate content for a presentation slide about:"
        try:
            content = self.llama_model.prompt(prompt)
            if self.presentation.current_slide:
                self.presentation.update_slide_content("Generated Slide", content)
                self.update_slide_view()
        except Exception as e:
            print(f"Error generating content: {e}")
            self.show_error_dialog("Failed to generate content")

    def on_research_clicked(self, widget):
        topic = "Artificial Intelligence in Office Suites"
        try:
            research = self.llama_model.prompt(f"Provide key points about: {topic}")
            if self.presentation.current_slide:
                self.presentation.update_slide_content(f"Research on {topic}", research)
                self.update_slide_view()
        except Exception as e:
            print(f"Error conducting research: {e}")
            self.show_error_dialog("Failed to conduct research")

    def show_error_dialog(self, message):
        dialog = Gtk.MessageDialog(
            transient_for=self,
            flags=0,
            message_type=Gtk.MessageType.ERROR,
            buttons=Gtk.ButtonsType.OK,
            text="Error",
        )
        dialog.format_secondary_text(message)
        dialog.run()
        dialog.destroy()

def main():
    app = PresentationUI()
    app.connect("destroy", Gtk.main_quit)
    app.show_all()
    Gtk.main()

if __name__ == "__main__":
    main()
